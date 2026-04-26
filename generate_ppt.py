from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hyper-Local Air Quality Forecasting"
    subtitle.text = "PM2.5 Intelligence & Health Alerts: Chennai-Ennore Corridor\nPresented by Dhanush Kumar T"

    # Slide 2: Problem Statement
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Problem Statement"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Current Challenges in Urban Air Quality:"
    p = tf.add_paragraph()
    p.text = "• Limited spatial resolution of monitoring stations."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Inability to provide street-level AQI forecasts."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Lack of personalized health alerts for sensitive populations."
    p.level = 1

    # Slide 3: Proposed Solution
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Proposed Solution: Multi-Agent System"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A decentralized architecture for intelligent forecasting:"
    p = tf.add_paragraph()
    p.text = "• Data Agent: Real-time multi-source data ingestion."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Spatial Agent: Advanced interpolation for street-level resolution."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Forecast Agent: Machine Learning models (XGBoost) for predictive trends."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Risk Agent: Persona-based health impact assessments."
    p.level = 1

    # Slide 4: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Architecture"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "End-to-End Pipeline:"
    p = tf.add_paragraph()
    p.text = "1. Data Collection: Weather, Industrial, and Traffic data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Calibration: Correcting sensor biases for accuracy."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Modeling: Comparing XGBoost vs. Linear Regression baselines."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Visualization: Interactive Folium heatmaps and Plotly charts."
    p.level = 1

    # Slide 5: Key Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Intelligent Dashboard Functionalities:"
    p = tf.add_paragraph()
    p.text = "• Street-Level Heatmaps: Visualizing AQI across Chennai."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Personalized Alerts: Tailored advice for Elderly, Children, and Outdoor Workers."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Dynamic Horizons: 6h, 12h, and 24h forecasting options."
    p.level = 1

    # Slide 6: Results
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Model Performance"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "XGBoost Implementation Results:"
    p = tf.add_paragraph()
    p.text = "• Outperformed traditional linear regression in accuracy."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Successful integration of multi-source features."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Real-time responsiveness for dashboard updates."
    p.level = 1

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion & Future Scope"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Next Steps:"
    p = tf.add_paragraph()
    p.text = "• Scaling to other cities across Tamil Nadu."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Integrating mobile app notifications for proximity-based alerts."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Incorporating additional pollutants (NO2, SO2, CO)."
    p.level = 1

    # Save the presentation
    output_path = "Air_Quality_Forecasting_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
